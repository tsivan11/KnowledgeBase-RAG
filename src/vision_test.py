"""
Standalone script to test GPT-4 Vision on embedded images from PDFs/PowerPoints.
Extracts images, analyzes them with Vision API, and reports results + costs.

Usage:
    python vision_test.py --file path/to/document.pdf
    python vision_test.py --file path/to/presentation.pptx
"""

import os
import base64
import argparse
from pathlib import Path
from io import BytesIO
from dotenv import load_dotenv

load_dotenv()

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    Presentation = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from openai import OpenAI
except ImportError:
    OpenAI = None

try:
    from PIL import Image
except ImportError:
    Image = None


def extract_images_from_pdf(pdf_path: Path):
    """Extract all embedded images from a PDF file."""
    if fitz is None:
        raise ImportError("PyMuPDF (fitz) not installed. Install with: pip install pymupdf")
    
    images = []
    doc = fitz.open(pdf_path)
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        image_list = page.get_images(full=True)
        
        for img_index, img_info in enumerate(image_list):
            xref = img_info[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            
            images.append({
                "page": page_num + 1,
                "index": img_index + 1,
                "bytes": image_bytes,
                "format": image_ext,
                "size": len(image_bytes)
            })
    
    doc.close()
    return images


def extract_images_from_pptx(pptx_path: Path):
    """Extract all embedded images from a PowerPoint file."""
    if Presentation is None:
        raise ImportError("python-pptx not installed. Install with: pip install python-pptx")
    
    images = []
    prs = Presentation(pptx_path)
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        img_index = 1
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_format = image.ext
                
                images.append({
                    "page": slide_num,
                    "index": img_index,
                    "bytes": image_bytes,
                    "format": image_format,
                    "size": len(image_bytes)
                })
                img_index += 1
    
    return images


def extract_images_from_docx(docx_path: Path):
    """Extract all embedded images from a Word document."""
    if Document is None:
        raise ImportError("python-docx not installed. Install with: pip install python-docx")
    
    images = []
    doc = Document(docx_path)
    
    # Extract from inline shapes (embedded images)
    img_index = 1
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                image_bytes = rel.target_part.blob
                image_format = rel.target_ref.split('.')[-1]
                
                images.append({
                    "page": None,  # Word doesn't have page numbers easily accessible
                    "index": img_index,
                    "bytes": image_bytes,
                    "format": image_format,
                    "size": len(image_bytes)
                })
                img_index += 1
            except:
                pass
    
    return images


def extract_images_from_xlsx(xlsx_path: Path):
    """Extract all embedded images from an Excel file."""
    if load_workbook is None:
        raise ImportError("openpyxl not installed. Install with: pip install openpyxl")
    
    images = []
    wb = load_workbook(xlsx_path)
    
    for sheet_num, sheet in enumerate(wb.worksheets, start=1):
        if hasattr(sheet, '_images') and sheet._images:
            for img_index, img in enumerate(sheet._images, start=1):
                try:
                    image_bytes = img._data()
                    # Try to get format from image
                    image_format = getattr(img, 'format', 'png')
                    
                    images.append({
                        "page": sheet_num,  # Using sheet number as page
                        "index": img_index,
                        "bytes": image_bytes,
                        "format": image_format,
                        "size": len(image_bytes)
                    })
                except Exception as e:
                    print(f"Warning: Could not extract image from sheet {sheet_num}: {e}")
    
    wb.close()
    return images


def analyze_image_with_vision(image_bytes: bytes, model: str = "gpt-4o-mini"):
    """Send image to GPT-4 Vision API for analysis."""
    if OpenAI is None:
        raise ImportError("OpenAI library not installed. Install with: pip install openai")
    
    if not os.getenv("OPENAI_API_KEY"):
        raise ValueError("OPENAI_API_KEY environment variable not set")
    
    client = OpenAI()
    
    # Encode image to base64
    base64_image = base64.b64encode(image_bytes).decode('utf-8')
    
    # Determine image format
    try:
        if Image:
            img = Image.open(BytesIO(image_bytes))
            format_str = img.format.lower() if img.format else "jpeg"
        else:
            format_str = "jpeg"
    except:
        format_str = "jpeg"
    
    # Call Vision API
    response = client.chat.completions.create(
        model=model,
        messages=[
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": "Describe this image in detail. Extract any text, describe any diagrams, charts, or visual elements. Be comprehensive."
                    },
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/{format_str};base64,{base64_image}",
                            "detail": "high"  # "low" or "high" - high gives better quality
                        }
                    }
                ]
            }
        ],
        max_tokens=1000
    )
    
    return {
        "description": response.choices[0].message.content,
        "input_tokens": response.usage.prompt_tokens,
        "output_tokens": response.usage.completion_tokens,
        "total_tokens": response.usage.total_tokens
    }


def estimate_cost(total_tokens: int, model: str = "gpt-4o-mini"):
    """Estimate cost based on token usage."""
    # Pricing as of Jan 2026 (approximate)
    prices = {
        "gpt-4o-mini": {"input": 0.15 / 1_000_000, "output": 0.60 / 1_000_000},
        "gpt-4o": {"input": 2.50 / 1_000_000, "output": 10.00 / 1_000_000}
    }
    
    if model not in prices:
        model = "gpt-4o-mini"
    
    # Rough estimate: assume 70% input, 30% output
    input_cost = total_tokens * 0.7 * prices[model]["input"]
    output_cost = total_tokens * 0.3 * prices[model]["output"]
    
    return input_cost + output_cost


def main():
    parser = argparse.ArgumentParser(description="Test GPT-4 Vision on embedded images from documents")
    parser.add_argument("--file", type=str, required=True, help="Path to PDF, PPTX, DOCX, or XLSX file")
    parser.add_argument("--model", type=str, default="gpt-4o-mini", 
                        choices=["gpt-4o-mini", "gpt-4o"],
                        help="Vision model to use (default: gpt-4o-mini)")
    parser.add_argument("--limit", type=int, default=None, 
                        help="Limit number of images to process (for testing)")
    args = parser.parse_args()
    
    file_path = Path(args.file)
    
    if not file_path.exists():
        print(f"Error: File not found: {file_path}")
        return
    
    # Extract images based on file type
    print(f"Processing: {file_path.name}")
    print(f"Using model: {args.model}\n")
    
    if file_path.suffix.lower() == ".pdf":
        images = extract_images_from_pdf(file_path)
    elif file_path.suffix.lower() in [".pptx", ".ppt"]:
        images = extract_images_from_pptx(file_path)
    elif file_path.suffix.lower() == ".docx":
        images = extract_images_from_docx(file_path)
    elif file_path.suffix.lower() in [".xlsx", ".xls"]:
        images = extract_images_from_xlsx(file_path)
    else:
        print(f"Error: Unsupported file type: {file_path.suffix}")
        print("Supported: .pdf, .pptx, .docx, .xlsx")
        return
    
    if not images:
        print("[ERROR] No embedded images found in document")
        return
    
    print(f"[OK] Found {len(images)} embedded image(s)\n")
    
    # Limit if specified
    if args.limit:
        images = images[:args.limit]
        print(f"[WARNING] Processing only first {args.limit} image(s)\n")
    
    # Process each image
    total_tokens = 0
    results = []
    
    for idx, img_info in enumerate(images, start=1):
        print(f"[Image {idx}/{len(images)}] Page/Slide {img_info['page']}, #{img_info['index']}")
        print(f"   Format: {img_info['format']}, Size: {img_info['size']:,} bytes")
        
        try:
            result = analyze_image_with_vision(img_info["bytes"], model=args.model)
            total_tokens += result["total_tokens"]
            
            print(f"   Tokens: {result['total_tokens']} (input: {result['input_tokens']}, output: {result['output_tokens']})")
            print(f"   Description: {result['description'][:150]}...")
            print()
            
            results.append({
                **img_info,
                **result
            })
            
        except Exception as e:
            print(f"   [ERROR] {e}")
            print()
    
    # Summary
    print("=" * 70)
    print("SUMMARY")
    print("=" * 70)
    print(f"Total images processed: {len(results)}")
    print(f"Total tokens used: {total_tokens:,}")
    estimated_cost = estimate_cost(total_tokens, args.model)
    print(f"Estimated cost: ${estimated_cost:.6f} (~${estimated_cost * 1000:.3f} per 1,000 images)")
    print()
    
    # Show full descriptions
    print("=" * 70)
    print("FULL DESCRIPTIONS")
    print("=" * 70)
    for idx, result in enumerate(results, start=1):
        print(f"\n[Image {idx} - Page/Slide {result['page']}, #{result['index']}]")
        print(result['description'])
        print("-" * 70)


if __name__ == "__main__":
    main()
