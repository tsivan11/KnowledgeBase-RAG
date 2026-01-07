# src/loaders.py
"""
Multi-format document loaders.
Each loader yields normalized records: {source, source_type, section, page, text}
"""
import json
import logging
import os
from pathlib import Path
from typing import Iterator, Dict, Any
from dotenv import load_dotenv

load_dotenv()

logger = logging.getLogger(__name__)

# Vision API configuration
USE_VISION_API = os.getenv("USE_VISION_API", "false").lower() == "true"
VISION_MODEL = os.getenv("VISION_MODEL", "gpt-4o-mini")
MAX_IMAGES_PER_FILE = int(os.getenv("MAX_IMAGES_PER_FILE", "5"))

# Optional dependencies
try:
    import fitz
except ImportError:
    fitz = None

try:
    from openai import OpenAI
except ImportError:
    OpenAI = None

try:
    import pytesseract
    from PIL import Image
except ImportError:
    pytesseract = None
    Image = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None


def _analyze_image_with_vision(image_bytes: bytes) -> str:
    """Analyze image using GPT-4 Vision API."""
    if not USE_VISION_API or OpenAI is None:
        return ""
    
    import base64
    from io import BytesIO
    
    try:
        client = OpenAI()
        base64_image = base64.b64encode(image_bytes).decode('utf-8')
        
        # Determine format
        try:
            if Image:
                img = Image.open(BytesIO(image_bytes))
                format_str = img.format.lower() if img.format else "jpeg"
            else:
                format_str = "jpeg"
        except:
            format_str = "jpeg"
        
        response = client.chat.completions.create(
            model=VISION_MODEL,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Describe this image in detail. Extract any text, describe any diagrams, charts, tables, or visual elements. Be comprehensive and focus on information that would be useful for answering questions about this document."
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/{format_str};base64,{base64_image}",
                                "detail": "high"
                            }
                        }
                    ]
                }
            ],
            max_tokens=1000
        )
        return response.choices[0].message.content
    except Exception as e:
        logger.debug(f"Vision API failed: {e}")
        return ""


def _format_table(table) -> str:
    """Format a pdfplumber table into a readable text representation."""
    if not table:
        return ""
    
    try:
        # Filter out None rows and cells
        rows = []
        for row in table:
            if row:
                cleaned_row = [str(cell).strip() if cell is not None else "" for cell in row]
                rows.append(cleaned_row)
        
        if not rows:
            return ""
        
        # Calculate column widths
        col_widths = [max(len(str(row[i])) for row in rows if i < len(row)) for i in range(len(rows[0]))]
        
        # Format rows
        formatted_rows = []
        for row in rows:
            formatted_row = " | ".join(str(cell).ljust(col_widths[i]) for i, cell in enumerate(row))
            formatted_rows.append(formatted_row)
        
        return "\n".join(formatted_rows)
    except Exception as e:
        logger.debug(f"Table formatting failed: {e}")
        return str(table)


def _ocr_page(page) -> str:
    """OCR a PDF page if pytesseract is available."""
    if pytesseract is None or Image is None:
        return ""
    
    import os
    tess_cmd = os.getenv("TESSERACT_CMD")
    if tess_cmd:
        pytesseract.pytesseract.tesseract_cmd = tess_cmd
    else:
        default_win = r"C:\\Program Files\\Tesseract-OCR\\tesseract.exe"
        if os.path.exists(default_win):
            pytesseract.pytesseract.tesseract_cmd = default_win
    
    try:
        pix = page.get_pixmap(dpi=200, alpha=False)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        text = pytesseract.image_to_string(img)
        return (text or "").strip()
    except Exception as e:
        logger.debug(f"OCR failed: {e}")
        return ""


def load_pdf(file_path: Path) -> Iterator[Dict[str, Any]]:
    """Load PDF with table extraction, OCR fallback for scanned pages."""
    if fitz is None:
        logger.warning(f"PyMuPDF not installed, skipping {file_path}")
        return
    
    source = file_path.name
    
    # Try table extraction with pdfplumber first
    if pdfplumber is not None:
        try:
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages, start=1):
                    page_text_parts = []
                    
                    # Extract tables
                    tables = page.extract_tables()
                    if tables:
                        for table_idx, table in enumerate(tables):
                            formatted_table = _format_table(table)
                            if formatted_table:
                                page_text_parts.append(f"[Table {table_idx + 1}]\n{formatted_table}")
                    
                    # Extract remaining text (non-table content)
                    regular_text = page.extract_text()
                    if regular_text:
                        page_text_parts.append(regular_text.strip())
                    
                    # Combine tables and text
                    combined_text = "\n\n".join(page_text_parts).strip()
                    
                    if combined_text:
                        yield {
                            "source": source,
                            "source_type": "pdf",
                            "page": i,
                            "section": None,
                            "text": combined_text,
                        }
            return  # Successfully processed with pdfplumber
        except Exception as e:
            logger.debug(f"pdfplumber extraction failed, falling back to PyMuPDF: {e}")
    
    # Fallback to PyMuPDF for regular text extraction
    try:
        doc = fitz.open(file_path)
        
        for i in range(len(doc)):
            page = doc.load_page(i)
            text = (page.get_text("text") or "").strip()
            
            if not text:
                text = _ocr_page(page)
            
            if not text:
                continue
            
            yield {
                "source": source,
                "source_type": "pdf",
                "page": i + 1,
                "section": None,
                "text": text,
            }
        
        doc.close()
    except Exception as e:
        logger.warning(f"Failed to load PDF {file_path}: {e}")


def load_txt(file_path: Path) -> Iterator[Dict[str, Any]]:
    """Load plain text file; treat each paragraph as a section."""
    try:
        content = file_path.read_text(encoding="utf-8")
        source = file_path.name
        
        # Split by double newlines (paragraphs)
        sections = [p.strip() for p in content.split("\n\n") if p.strip()]
        
        for idx, text in enumerate(sections, start=1):
            yield {
                "source": source,
                "source_type": "txt",
                "page": None,
                "section": idx,
                "text": text,
            }
    except Exception as e:
        logger.warning(f"Failed to load TXT {file_path}: {e}")


def load_md(file_path: Path) -> Iterator[Dict[str, Any]]:
    """Load Markdown file; preserve headings as section metadata."""
    try:
        content = file_path.read_text(encoding="utf-8")
        source = file_path.name
        
        # Simple approach: split by headings (# ## ###)
        import re
        parts = re.split(r"^(#{1,6}\s+.+)$", content, flags=re.MULTILINE)
        
        section_idx = 1
        current_heading = None
        
        for part in parts:
            if not part.strip():
                continue
            
            # Check if it's a heading
            heading_match = re.match(r"^(#{1,6})\s+(.+)$", part)
            if heading_match:
                current_heading = heading_match.group(2).strip()
            else:
                text = part.strip()
                if text:
                    yield {
                        "source": source,
                        "source_type": "md",
                        "page": None,
                        "section": current_heading or f"section_{section_idx}",
                        "text": text,
                    }
                    section_idx += 1
    except Exception as e:
        logger.warning(f"Failed to load MD {file_path}: {e}")


def load_docx(file_path: Path) -> Iterator[Dict[str, Any]]:
    """Load DOCX file; extract paragraphs and images."""
    if Document is None:
        logger.warning(f"python-docx not installed, skipping {file_path}")
        return
    
    try:
        doc = Document(file_path)
        source = file_path.name
        
        # Extract text paragraphs
        text_parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                text_parts.append(text)
        
        # Extract images if vision enabled
        image_descriptions = []
        total_images = 0
        if USE_VISION_API:
            img_index = 1
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    total_images += 1
                    if MAX_IMAGES_PER_FILE > 0 and img_index > MAX_IMAGES_PER_FILE:
                        continue
                    try:
                        image_bytes = rel.target_part.blob
                        description = _analyze_image_with_vision(image_bytes)
                        if description:
                            image_descriptions.append(f"[Image {img_index}]: {description}")
                            img_index += 1
                    except:
                        pass
            if MAX_IMAGES_PER_FILE > 0 and total_images > MAX_IMAGES_PER_FILE:
                logger.warning(f"Processed {MAX_IMAGES_PER_FILE}/{total_images} images in {source} (limit reached)")
        
        # Combine text and images
        content_parts = []
        if text_parts:
            content_parts.append("\n\n".join(text_parts))
        if image_descriptions:
            content_parts.append("\n\n=== EMBEDDED IMAGES ===\n\n" + "\n\n".join(image_descriptions))
        
        combined_text = "\n\n".join(content_parts).strip()
        if combined_text:
            yield {
                "source": source,
                "source_type": "docx",
                "page": None,
                "section": None,
                "text": combined_text,
            }
    except Exception as e:
        logger.warning(f"Failed to load DOCX {file_path}: {e}")


def load_html(file_path: Path) -> Iterator[Dict[str, Any]]:
    """Load HTML file; extract text and preserve heading hierarchy."""
    if BeautifulSoup is None:
        logger.warning(f"BeautifulSoup4 not installed, skipping {file_path}")
        return
    
    try:
        content = file_path.read_text(encoding="utf-8")
        soup = BeautifulSoup(content, "html.parser")
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        source = file_path.name
        section_idx = 1
        current_heading = None
        
        # Walk through body or full soup
        for elem in (soup.body or soup).descendants:
            if isinstance(elem, str):
                text = elem.strip()
                if text and len(text) > 1:  # Avoid single chars
                    yield {
                        "source": source,
                        "source_type": "html",
                        "page": None,
                        "section": current_heading or f"section_{section_idx}",
                        "text": text,
                    }
                    section_idx += 1
            elif hasattr(elem, "name"):
                if elem.name in ["h1", "h2", "h3"]:
                    current_heading = elem.get_text().strip()
    except Exception as e:
        logger.warning(f"Failed to load HTML {file_path}: {e}")


def load_csv(file_path: Path) -> Iterator[Dict[str, Any]]:
    """Load CSV file; flatten rows with column context."""
    if pd is None:
        logger.warning(f"pandas not installed, skipping {file_path}")
        return
    
    try:
        df = pd.read_csv(file_path)
        source = file_path.name
        
        for idx, row in df.iterrows():
            # Format row as key=value pairs
            text_parts = []
            for col, val in row.items():
                if pd.notna(val):
                    text_parts.append(f"{col}: {val}")
            
            text = " | ".join(text_parts)
            if text.strip():
                yield {
                    "source": source,
                    "source_type": "csv",
                    "page": None,
                    "section": f"row_{idx + 2}",  # +2 for header and 0-indexing
                    "text": text,
                }
    except Exception as e:
        logger.warning(f"Failed to load CSV {file_path}: {e}")


def load_xlsx(file_path: Path) -> Iterator[Dict[str, Any]]:
    """Load Excel file; process each sheet with text and images."""
    if pd is None:
        logger.warning(f"pandas not installed, skipping {file_path}")
        return
    
    try:
        # Read text data with pandas
        excel_file = pd.ExcelFile(file_path)
        source = file_path.name
        
        # Extract images if vision enabled (requires openpyxl)
        sheet_images = {}
        total_images_processed = 0
        if USE_VISION_API:
            try:
                from openpyxl import load_workbook
                wb = load_workbook(file_path)
                for sheet in wb.worksheets:
                    images = []
                    if hasattr(sheet, '_images') and sheet._images:
                        for img_index, img in enumerate(sheet._images, start=1):
                            if MAX_IMAGES_PER_FILE > 0 and total_images_processed >= MAX_IMAGES_PER_FILE:
                                break
                            try:
                                image_bytes = img._data()
                                description = _analyze_image_with_vision(image_bytes)
                                if description:
                                    images.append(f"[Image {img_index}]: {description}")
                                    total_images_processed += 1
                            except:
                                pass
                    if images:
                        sheet_images[sheet.title] = images
                    if MAX_IMAGES_PER_FILE > 0 and total_images_processed >= MAX_IMAGES_PER_FILE:
                        break
                wb.close()
            except ImportError:
                pass  # openpyxl not available
            except Exception as e:
                logger.debug(f"Failed to extract Excel images: {e}")
        
        # Process each sheet
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            
            # Build text from rows
            text_parts = []
            for idx, row in df.iterrows():
                row_parts = []
                for col, val in row.items():
                    if pd.notna(val):
                        row_parts.append(f"{col}: {val}")
                if row_parts:
                    text_parts.append(" | ".join(row_parts))
            
            # Combine text and images for this sheet
            content_parts = []
            if text_parts:
                content_parts.append("\n".join(text_parts))
            if sheet_name in sheet_images:
                content_parts.append("\n\n=== EMBEDDED IMAGES ===\n\n" + "\n\n".join(sheet_images[sheet_name]))
            
            combined_text = "\n\n".join(content_parts).strip()
            if combined_text:
                yield {
                    "source": source,
                    "source_type": "xlsx",
                    "page": None,
                    "section": sheet_name,
                    "text": combined_text,
                }
    except Exception as e:
        logger.warning(f"Failed to load Excel {file_path}: {e}")


def load_pptx(file_path: Path) -> Iterator[Dict[str, Any]]:
    """Load PowerPoint file; extract text and images from slides."""
    if Presentation is None:
        logger.warning(f"python-pptx not installed, skipping {file_path}")
        return
    
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
        source = file_path.name
        total_images_in_file = 0
        
        for slide_idx, slide in enumerate(prs.slides, start=1):
            # Extract text from all shapes
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text.strip())
            
            # Extract images if vision enabled
            image_descriptions = []
            if USE_VISION_API:
                img_index = 1
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        if MAX_IMAGES_PER_FILE > 0 and total_images_in_file >= MAX_IMAGES_PER_FILE:
                            break
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            description = _analyze_image_with_vision(image_bytes)
                            if description:
                                image_descriptions.append(f"[Image {img_index}]: {description}")
                                img_index += 1
                                total_images_in_file += 1
                        except:
                            pass
            
            # Combine text and images
            content_parts = []
            if text_parts:
                content_parts.append("\n".join(text_parts))
            if image_descriptions:
                content_parts.append("\n\n=== EMBEDDED IMAGES ===\n\n" + "\n\n".join(image_descriptions))
            
            combined_text = "\n\n".join(content_parts).strip()
            if combined_text:
                yield {
                    "source": source,
                    "source_type": "pptx",
                    "page": slide_idx,
                    "section": None,
                    "text": combined_text,
                }
    except Exception as e:
        logger.warning(f"Failed to load PowerPoint {file_path}: {e}")


def load_image(file_path: Path) -> Iterator[Dict[str, Any]]:
    """Load image file and extract text via Vision API or OCR."""
    try:
        image_bytes = file_path.read_bytes()
        text = ""
        
        # Try Vision API first if enabled
        if USE_VISION_API:
            text = _analyze_image_with_vision(image_bytes)
            if text:
                logger.debug(f"Used Vision API for {file_path.name}")
        
        # Fallback to Tesseract OCR if vision disabled or failed
        if not text and pytesseract is not None and Image is not None:
            import os
            tess_cmd = os.getenv("TESSERACT_CMD")
            if tess_cmd:
                pytesseract.pytesseract.tesseract_cmd = tess_cmd
            else:
                default_win = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
                if os.path.exists(default_win):
                    pytesseract.pytesseract.tesseract_cmd = default_win
            
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img)
            text = (text or "").strip()
            if text:
                logger.debug(f"Used Tesseract OCR for {file_path.name}")
        
        if text:
            yield {
                "source": file_path.name,
                "source_type": file_path.suffix.lower()[1:],
                "page": None,
                "section": None,
                "text": text,
            }
    except Exception as e:
        logger.warning(f"Failed to load image {file_path}: {e}")


def load_audio(file_path: Path) -> Iterator[Dict[str, Any]]:
    """Transcribe audio file using OpenAI Whisper API."""
    if OpenAI is None:
        logger.warning(f"OpenAI not installed, skipping {file_path}")
        return
    
    import os
    if not os.getenv("OPENAI_API_KEY"):
        logger.warning(f"OPENAI_API_KEY not set, skipping {file_path}")
        return
    
    try:
        client = OpenAI()
        
        # Check file size (Whisper API limit is 25MB)
        file_size_mb = file_path.stat().st_size / (1024 * 1024)
        if file_size_mb > 25:
            logger.warning(f"Audio file {file_path.name} is {file_size_mb:.1f}MB, exceeds 25MB limit")
            return
        
        logger.info(f"Transcribing audio file: {file_path.name} ({file_size_mb:.1f}MB)...")
        
        with open(file_path, "rb") as audio_file:
            transcript = client.audio.transcriptions.create(
                model="whisper-1",
                file=audio_file,
                response_format="text"
            )
        
        # Return full transcript as single page
        # Let the chunker handle splitting intelligently
        if transcript and transcript.strip():
            yield {
                "source": file_path.name,
                "source_type": "audio",
                "page": None,
                "section": None,
                "text": transcript.strip(),
            }
        
        logger.info(f"Successfully transcribed {file_path.name}")
        
    except Exception as e:
        logger.warning(f"Failed to transcribe audio {file_path}: {e}")


# Loader registry
LOADERS = {
    ".pdf": load_pdf,
    ".txt": load_txt,
    ".md": load_md,
    ".docx": load_docx,
    ".html": load_html,
    ".htm": load_html,
    ".csv": load_csv,
    ".xlsx": load_xlsx,
    ".xls": load_xlsx,
    ".pptx": load_pptx,
    ".jpg": load_image,
    ".jpeg": load_image,
    ".png": load_image,
    ".tiff": load_image,
    ".tif": load_image,
    ".bmp": load_image,
    # Audio formats
    ".mp3": load_audio,
    ".mp4": load_audio,
    ".mpeg": load_audio,
    ".mpga": load_audio,
    ".m4a": load_audio,
    ".wav": load_audio,
    ".webm": load_audio,
}


def load_file(file_path: Path) -> Iterator[Dict[str, Any]]:
    """Auto-detect format and load file."""
    ext = file_path.suffix.lower()
    
    if ext not in LOADERS:
        logger.warning(f"Unsupported file type: {ext} ({file_path.name})")
        return
    
    loader = LOADERS[ext]
    yield from loader(file_path)
