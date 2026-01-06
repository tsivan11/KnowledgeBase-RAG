"""
Ingest documents with embedded image extraction using GPT-4 Vision.
Creates pages.jsonl compatible with existing chunk_pages.py and build_index.py pipeline.

Usage:
    python src/vision_ingest.py --domain test
    python src/vision_ingest.py --domain test --file kb/test/document.pdf
    python src/vision_ingest.py --domain test --file kb/test/presentation.pptx
    python src/vision_ingest.py --domain test --file kb/test/report.docx
    python src/vision_ingest.py --domain test --file kb/test/data.xlsx
"""

import os
import json
import base64
import logging
import argparse
from pathlib import Path
from io import BytesIO
from dotenv import load_dotenv

load_dotenv()

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    Presentation = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from openai import OpenAI
except ImportError:
    OpenAI = None

try:
    from PIL import Image
except ImportError:
    Image = None

logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")
logger = logging.getLogger(__name__)


def get_paths(domain: str):
    """Get input and output paths for domain."""
    if Path.cwd().name == "src":
        base = Path("..")
    else:
        base = Path(".")
    
    input_dir = base / "kb" / domain
    out_path = base / "data" / domain / "pages.jsonl"
    return input_dir, out_path


def analyze_image_with_vision(image_bytes: bytes, model: str = "gpt-4o-mini") -> str:
    """Send image to GPT-4 Vision API for analysis."""
    if OpenAI is None:
        raise ImportError("OpenAI library not installed")
    
    if not os.getenv("OPENAI_API_KEY"):
        raise ValueError("OPENAI_API_KEY environment variable not set")
    
    client = OpenAI()
    
    # Encode image to base64
    base64_image = base64.b64encode(image_bytes).decode('utf-8')
    
    # Determine image format
    try:
        if Image:
            img = Image.open(BytesIO(image_bytes))
            format_str = img.format.lower() if img.format else "jpeg"
        else:
            format_str = "jpeg"
    except:
        format_str = "jpeg"
    
    # Call Vision API
    response = client.chat.completions.create(
        model=model,
        messages=[
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": "Describe this image in detail. Extract any text, describe any diagrams, charts, tables, or visual elements. Be comprehensive and focus on information that would be useful for answering questions about this document."
                    },
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/{format_str};base64,{base64_image}",
                            "detail": "high"
                        }
                    }
                ]
            }
        ],
        max_tokens=1000
    )
    
    return response.choices[0].message.content


def extract_pdf_with_vision(file_path: Path, model: str = "gpt-4o-mini"):
    """
    Extract text and images from PDF.
    Returns list of page records with text and embedded image descriptions.
    """
    if fitz is None:
        raise ImportError("PyMuPDF not installed")
    
    records = []
    doc = fitz.open(file_path)
    source = file_path.name
    
    logger.info(f"Processing {source} ({len(doc)} pages)")
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        page_number = page_num + 1
        
        # Extract text
        text = (page.get_text("text") or "").strip()
        
        # Extract embedded images
        image_list = page.get_images(full=True)
        image_descriptions = []
        
        if image_list:
            logger.info(f"  Page {page_number}: Found {len(image_list)} embedded image(s)")
            
            for img_index, img_info in enumerate(image_list, start=1):
                try:
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    logger.info(f"    Analyzing image {img_index}/{len(image_list)} with vision API...")
                    description = analyze_image_with_vision(image_bytes, model)
                    image_descriptions.append(f"[Image {img_index}]: {description}")
                    
                except Exception as e:
                    logger.warning(f"    Failed to analyze image {img_index}: {e}")
        
        # Combine text and image descriptions
        content_parts = []
        if text:
            content_parts.append(text)
        if image_descriptions:
            content_parts.append("\n\n=== EMBEDDED IMAGES ===\n\n" + "\n\n".join(image_descriptions))
        
        combined_text = "\n\n".join(content_parts).strip()
        
        if combined_text:
            records.append({
                "source": source,
                "source_type": "pdf",
                "page": page_number,
                "section": None,
                "text": combined_text,
            })
            logger.info(f"  Page {page_number}: {len(text)} chars text + {len(image_descriptions)} images")
    
    doc.close()
    return records


def extract_pptx_with_vision(file_path: Path, model: str = "gpt-4o-mini"):
    """
    Extract text and images from PowerPoint.
    Returns list of slide records with text and embedded image descriptions.
    """
    if Presentation is None:
        raise ImportError("python-pptx not installed")
    
    records = []
    prs = Presentation(file_path)
    source = file_path.name
    
    logger.info(f"Processing {source} ({len(prs.slides)} slides)")
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        # Extract text from slide
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text.strip())
        
        text = "\n".join(text_parts).strip()
        
        # Extract images from slide
        image_descriptions = []
        img_index = 1
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    image_bytes = image.blob
                    
                    logger.info(f"  Slide {slide_num}: Analyzing image {img_index} with vision API...")
                    description = analyze_image_with_vision(image_bytes, model)
                    image_descriptions.append(f"[Image {img_index}]: {description}")
                    img_index += 1
                    
                except Exception as e:
                    logger.warning(f"    Failed to analyze image {img_index}: {e}")
                    img_index += 1
        
        if image_descriptions:
            logger.info(f"  Slide {slide_num}: Found {len(image_descriptions)} image(s)")
        
        # Combine text and image descriptions
        content_parts = []
        if text:
            content_parts.append(text)
        if image_descriptions:
            content_parts.append("\n\n=== EMBEDDED IMAGES ===\n\n" + "\n\n".join(image_descriptions))
        
        combined_text = "\n\n".join(content_parts).strip()
        
        if combined_text:
            records.append({
                "source": source,
                "source_type": "pptx",
                "page": slide_num,  # Using page number for slide number
                "section": None,
                "text": combined_text,
            })
            logger.info(f"  Slide {slide_num}: {len(text)} chars text + {len(image_descriptions)} images")
    
    return records


def extract_docx_with_vision(file_path: Path, model: str = "gpt-4o-mini"):
    """Extract text and images from Word document."""
    if Document is None:
        raise ImportError("python-docx not installed")
    
    records = []
    doc = Document(file_path)
    source = file_path.name
    
    logger.info(f"Processing {source}")
    
    # Extract text from paragraphs
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())
    
    text = "\n\n".join(text_parts).strip()
    
    # Extract images
    image_descriptions = []
    img_index = 1
    
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                image_bytes = rel.target_part.blob
                logger.info(f"  Analyzing image {img_index} with vision API...")
                description = analyze_image_with_vision(image_bytes, model)
                image_descriptions.append(f"[Image {img_index}]: {description}")
                img_index += 1
            except Exception as e:
                logger.warning(f"    Failed to analyze image {img_index}: {e}")
                img_index += 1
    
    if image_descriptions:
        logger.info(f"  Found {len(image_descriptions)} image(s)")
    
    # Combine text and images
    content_parts = []
    if text:
        content_parts.append(text)
    if image_descriptions:
        content_parts.append("\n\n=== EMBEDDED IMAGES ===\n\n" + "\n\n".join(image_descriptions))
    
    combined_text = "\n\n".join(content_parts).strip()
    
    if combined_text:
        records.append({
            "source": source,
            "source_type": "docx",
            "page": None,
            "section": None,
            "text": combined_text,
        })
        logger.info(f"  {len(text)} chars text + {len(image_descriptions)} images")
    
    return records


def extract_xlsx_with_vision(file_path: Path, model: str = "gpt-4o-mini"):
    """Extract text and images from Excel file."""
    if load_workbook is None:
        raise ImportError("openpyxl not installed")
    
    records = []
    wb = load_workbook(file_path, data_only=True)
    source = file_path.name
    
    logger.info(f"Processing {source} ({len(wb.worksheets)} sheets)")
    
    for sheet_num, sheet in enumerate(wb.worksheets, start=1):
        # Extract cell values
        text_parts = []
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
            if row_text.strip():
                text_parts.append(row_text)
        
        text = "\n".join(text_parts).strip()
        
        # Extract images from sheet
        image_descriptions = []
        if hasattr(sheet, '_images') and sheet._images:
            logger.info(f"  Sheet {sheet_num} ({sheet.title}): Found {len(sheet._images)} image(s)")
            
            for img_index, img in enumerate(sheet._images, start=1):
                try:
                    image_bytes = img._data()
                    logger.info(f"    Analyzing image {img_index} with vision API...")
                    description = analyze_image_with_vision(image_bytes, model)
                    image_descriptions.append(f"[Image {img_index}]: {description}")
                except Exception as e:
                    logger.warning(f"    Failed to analyze image {img_index}: {e}")
        
        # Combine text and images for this sheet
        content_parts = []
        if text:
            content_parts.append(f"Sheet: {sheet.title}\n\n{text}")
        if image_descriptions:
            content_parts.append("\n\n=== EMBEDDED IMAGES ===\n\n" + "\n\n".join(image_descriptions))
        
        combined_text = "\n\n".join(content_parts).strip()
        
        if combined_text:
            records.append({
                "source": source,
                "source_type": "xlsx",
                "page": sheet_num,
                "section": sheet.title,
                "text": combined_text,
            })
            logger.info(f"  Sheet {sheet_num}: {len(text)} chars text + {len(image_descriptions)} images")
    
    wb.close()
    return records


def ingest_with_vision(domain: str, specific_file: str = None, model: str = "gpt-4o-mini"):
    """Ingest documents with vision API, write to pages.jsonl."""
    input_dir, out_path = get_paths(domain)
    
    if not input_dir.exists():
        raise FileNotFoundError(f"Missing folder: {input_dir.resolve()}")
    
    out_path.parent.mkdir(parents=True, exist_ok=True)
    
    # Get files to process
    if specific_file:
        files = [Path(specific_file)]
        if not files[0].exists():
            raise FileNotFoundError(f"File not found: {specific_file}")
    else:
        # Find all supported document types
        pdf_files = list(input_dir.glob("**/*.pdf"))
        pptx_files = list(input_dir.glob("**/*.pptx"))
        docx_files = list(input_dir.glob("**/*.docx"))
        xlsx_files = list(input_dir.glob("**/*.xlsx"))
        files = sorted(pdf_files + pptx_files + docx_files + xlsx_files)
    
    if not files:
        raise FileNotFoundError(f"No supported files found (PDF, PPTX, DOCX, XLSX)")
    
    total_records = 0
    logger.info(f"Processing {len(files)} file(s) with vision API (model: {model})")
    
    with out_path.open("w", encoding="utf-8") as f:
        for file_path in files:
            try:
                # Process based on file type
                if file_path.suffix.lower() == ".pdf":
                    records = extract_pdf_with_vision(file_path, model)
                elif file_path.suffix.lower() in [".pptx", ".ppt"]:
                    records = extract_pptx_with_vision(file_path, model)
                elif file_path.suffix.lower() == ".docx":
                    records = extract_docx_with_vision(file_path, model)
                elif file_path.suffix.lower() in [".xlsx", ".xls"]:
                    records = extract_xlsx_with_vision(file_path, model)
                else:
                    logger.warning(f"Unsupported file type: {file_path.suffix}")
                    continue
                
                for rec in records:
                    f.write(json.dumps(rec, ensure_ascii=False) + "\n")
                    total_records += 1
                
            except Exception as e:
                logger.warning(f"{file_path.name}: FAILED - {e}")
    
    logger.info(f"Wrote {total_records} page record(s) to {out_path}")
    logger.info(f"\nNext steps:")
    logger.info(f"  1. Chunk: python src/chunk_pages.py --domain {domain}")
    logger.info(f"  2. Index: python src/build_index.py --domain {domain}")
    logger.info(f"  3. Query: python src/ask.py --domain {domain} --query 'your question'")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Ingest documents (PDF/PPTX/DOCX/XLSX) with GPT-4 Vision for images")
    parser.add_argument("--domain", type=str, required=True,
                        help="Domain folder name (e.g., 'test')")
    parser.add_argument("--file", type=str, default=None,
                        help="Process specific file instead of entire domain folder")
    parser.add_argument("--model", type=str, default="gpt-4o-mini",
                        choices=["gpt-4o-mini", "gpt-4o"],
                        help="Vision model to use (default: gpt-4o-mini)")
    args = parser.parse_args()
    
    ingest_with_vision(args.domain, args.file, args.model)
